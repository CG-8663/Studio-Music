#!/usr/bin/env python3
"""Generate Major 72 Studio Investor Pitch Deck"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
BLACK = RGBColor(0x12, 0x12, 0x12)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xD4, 0xA5, 0x37)
DARK_GRAY = RGBColor(0x2A, 0x2A, 0x2A)
MED_GRAY = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
ACCENT_BLUE = RGBColor(0x4A, 0x90, 0xD9)
RED_ACCENT = RGBColor(0xE8, 0x4D, 0x4D)
GREEN_ACCENT = RGBColor(0x4D, 0xC9, 0x6D)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=BLACK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_paragraph(tf, text, font_size=18, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = alignment
    p.space_before = space_before
    return p

def add_shape_with_text(slide, left, top, width, height, text,
                        fill_color=DARK_GRAY, text_color=WHITE,
                        font_size=14, bold=False, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(shape, Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.font.name = "Calibri"
    return tf

def add_gold_line(slide, top):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(top),
                                   Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

add_text_box(slide, 1, 1.5, 11, 1.2, "MAJOR 72 STUDIO", 54, GOLD, True, PP_ALIGN.CENTER)
add_gold_line(slide, 2.8)
add_text_box(slide, 1, 3.2, 11, 0.8, "Your Music. Eddie's Ears.", 32, WHITE, False, PP_ALIGN.CENTER)
add_text_box(slide, 2, 4.3, 9, 0.8,
    "The Professional-AI Hybrid Music Production Platform",
    20, MED_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 4, 6.0, 5, 0.5, "Investor Presentation  |  April 2026",
    14, MED_GRAY, False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: THE PROBLEM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.5, 11, 0.8, "THE PROBLEM", 36, GOLD, True)
add_gold_line(slide, 1.2)

add_text_box(slide, 0.8, 1.6, 11, 0.8,
    "AI music is exploding. Professionals don't trust it.",
    26, WHITE, True)

# Three problem columns
problems = [
    ("50,000", "AI tracks flood\nstreaming platforms\nEVERY DAY",
     "Spotify removed 75M\n\"spammy tracks\" last year"),
    ("80%", "of producers are\nneutral-to-negative\nabout AI tools",
     "Only 20% describe their\nfeelings as positive"),
    ("$0", "copyright value of\npurely AI-generated\nmusic",
     "USCO: prompts alone don't\ncreate copyrightable works"),
]

for i, (big_num, desc, sub) in enumerate(problems):
    x = 0.8 + (i * 4.1)
    add_shape_with_text(slide, x, 2.8, 3.7, 3.8, "", DARK_GRAY)
    add_text_box(slide, x + 0.3, 2.9, 3.1, 0.8, big_num, 48, RED_ACCENT, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.3, 3.8, 3.1, 1.2, desc, 16, WHITE, False, PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.3, 5.2, 3.1, 1.0, sub, 12, MED_GRAY, False, PP_ALIGN.CENTER)

add_text_box(slide, 0.8, 6.8, 11, 0.5,
    "The gap: professional creative judgment — knowing how to make a mix breathe — is the one thing AI can't generate alone.",
    14, GOLD, True, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 3: THE MARKET
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.5, 11, 0.8, "THE MARKET", 36, GOLD, True)
add_gold_line(slide, 1.2)

# Market size boxes
metrics = [
    ("$2.4B", "Market Size (2024)", "Generative AI in Music"),
    ("$22.7B", "Projected (2035)", "22-30% CAGR"),
    ("$300M", "Suno ARR (Feb 2026)", "2M paid subscribers"),
    ("$75M", "Songfinch Revenue", "Custom songs (2023)"),
]

for i, (value, label, sub) in enumerate(metrics):
    x = 0.8 + (i * 3.1)
    add_shape_with_text(slide, x, 1.7, 2.7, 1.8, "", DARK_GRAY)
    add_text_box(slide, x + 0.2, 1.8, 2.3, 0.7, value, 36, GOLD, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, 2.5, 2.3, 0.4, label, 14, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, 2.9, 2.3, 0.4, sub, 12, MED_GRAY, False, PP_ALIGN.CENTER)

# Key stats
stats = [
    "87% of artists have used AI in their workflow (LANDR, 1,200+ artists)",
    "~50% of producers are \"cautious experimenters\" — interested but uncommitted",
    ">33% fear \"musical sameness\" — generic output that erodes their identity",
    "$591M+ in music tech startup funding in 2025 alone",
    "112,000+ tracks uploaded to streaming platforms daily — nearly half AI-touched",
]

tf = add_text_box(slide, 0.8, 4.0, 11.5, 3.0, "", 16, WHITE)
for stat in stats:
    add_paragraph(tf, f"  \u2022  {stat}", 15, LIGHT_GRAY, False, PP_ALIGN.LEFT, Pt(10))

add_text_box(slide, 0.8, 6.8, 11, 0.5,
    "Sources: Spherical Insights, MRFR, Grand View Research, TechCrunch, Billboard, LANDR, Sonarworks",
    10, MED_GRAY, False, PP_ALIGN.LEFT)

# ============================================================
# SLIDE 4: THE INSIGHT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.5, 11, 0.8, "THE INSIGHT", 36, GOLD, True)
add_gold_line(slide, 1.2)

add_text_box(slide, 1.5, 2.0, 10, 1.5,
    "\"Producers want AI for technical drudgery.\nThey want to keep the creative decisions.\"",
    30, WHITE, True, PP_ALIGN.CENTER)

add_text_box(slide, 1.5, 3.8, 10, 0.6,
    "— Sonarworks / Sound On Sound 2026 Survey (1,100+ working producers)",
    16, MED_GRAY, False, PP_ALIGN.CENTER)

# Two columns: what producers want AI for vs what they don't
add_shape_with_text(slide, 1.0, 4.6, 5.2, 2.2, "", DARK_GRAY)
tf = add_text_box(slide, 1.3, 4.7, 4.6, 0.4, "WANT AI FOR (Technical)", 16, GREEN_ACCENT, True)
for item in ["Stem separation (73.9%)", "Vocal tuning & drum editing",
             "File management & cleanup", "Reference mastering"]:
    add_paragraph(tf, f"  \u2713  {item}", 14, LIGHT_GRAY)

add_shape_with_text(slide, 7.0, 4.6, 5.2, 2.2, "", DARK_GRAY)
tf = add_text_box(slide, 7.3, 4.7, 4.6, 0.4, "DON'T WANT AI FOR (Creative)", 16, RED_ACCENT, True)
for item in ["Musical direction & taste", "Emotional judgment",
             "Arrangement decisions", "Artist identity & sound"]:
    add_paragraph(tf, f"  \u2717  {item}", 14, LIGHT_GRAY)

# ============================================================
# SLIDE 5: THE SOLUTION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.5, 11, 0.8, "THE SOLUTION", 36, GOLD, True)
add_gold_line(slide, 1.2)

add_text_box(slide, 0.8, 1.6, 11, 0.8,
    "Major 72 Studio encodes professional taste into AI-guided production.",
    22, WHITE, True)

# Flow diagram
steps = [
    ("Upload\nYour Idea", "Phone recording,\nbass riff, vocal demo"),
    ("Eddie's\nCreative Mode", "Professional judgment\nguides the AI"),
    ("AI Handles\nThe Work", "Separation, mixing,\nmastering in <60s"),
    ("Download\nPro Output", "Polished mix, stems,\nsession metadata"),
]

for i, (title, sub) in enumerate(steps):
    x = 0.6 + (i * 3.2)
    add_shape_with_text(slide, x, 2.7, 2.6, 1.5, title, DARK_GRAY, GOLD, 18, True)
    add_text_box(slide, x, 4.3, 2.6, 0.8, sub, 12, MED_GRAY, False, PP_ALIGN.CENTER)
    if i < 3:
        add_text_box(slide, x + 2.6, 3.2, 0.6, 0.5, "\u2192", 28, GOLD, True, PP_ALIGN.CENTER)

add_text_box(slide, 0.8, 5.4, 11, 1.5,
    "Key principle: the AI doesn't change your music. It doesn't add instruments.\n"
    "It takes what YOU created and makes it sound like a professional spent\n"
    "six hours on it. Your music. Eddie's ears. Nothing added. Everything elevated.",
    16, WHITE, False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 6: SERVICE MODEL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.5, 11, 0.8, "FOUR-TIER SERVICE MODEL", 36, GOLD, True)
add_gold_line(slide, 1.2)

tiers = [
    ("HEAR IT", "$29-49/mo", "Upload rough idea\nGet it back polished", "MVP", GOLD),
    ("BUILD IT", "$79-149/track", "Full AI arrangement\nwith Session Players", "Phase 2", ACCENT_BLUE),
    ("PLAY IT", "$299-999/project", "Stage-ready charts\n+ session musicians", "Phase 3", ACCENT_BLUE),
    ("GIFT IT", "$49-149/song", "Custom song from\nyour story", "MVP", GOLD),
]

for i, (name, price, desc, phase, accent) in enumerate(tiers):
    x = 0.5 + (i * 3.15)
    add_shape_with_text(slide, x, 1.7, 2.9, 3.5, "", DARK_GRAY)
    add_text_box(slide, x + 0.2, 1.8, 2.5, 0.5, name, 24, accent, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, 2.3, 2.5, 0.5, price, 18, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, 3.0, 2.5, 1.0, desc, 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    phase_color = GREEN_ACCENT if phase == "MVP" else MED_GRAY
    add_text_box(slide, x + 0.2, 4.3, 2.5, 0.4, phase, 13, phase_color, True, PP_ALIGN.CENTER)

# Land and expand
add_text_box(slide, 0.8, 5.5, 11, 0.5,
    "LAND AND EXPAND: Each tier validates the next. Don't build ahead of demand.",
    16, GOLD, True, PP_ALIGN.CENTER)

tf = add_text_box(slide, 1.5, 6.1, 10, 1.2, "", 14, LIGHT_GRAY)
add_paragraph(tf, "Hear It users ask \"I need drums and keys\"  \u2192  Build It", 14, LIGHT_GRAY)
add_paragraph(tf, "Build It users ask \"My band wants to play this live\"  \u2192  Play It", 14, LIGHT_GRAY)
add_paragraph(tf, "Gift It runs as a separate brand — same engine, different front door", 14, LIGHT_GRAY)

# ============================================================
# SLIDE 7: COMPETITIVE POSITIONING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.5, 11, 0.8, "THE EMPTY QUADRANT", 36, GOLD, True)
add_gold_line(slide, 1.2)

add_text_box(slide, 0.8, 1.5, 11, 0.6,
    "No competitor occupies \"Professional Quality + Human-Guided AI\"",
    20, WHITE, True, PP_ALIGN.CENTER)

# Positioning map as text
add_shape_with_text(slide, 1.5, 2.3, 10.3, 4.5, "", DARK_GRAY)

# Quadrant labels
add_text_box(slide, 5.5, 2.4, 3, 0.4, "PROFESSIONAL QUALITY", 12, MED_GRAY, True, PP_ALIGN.CENTER)
add_text_box(slide, 5.5, 6.3, 3, 0.4, "CONSUMER / CASUAL", 12, MED_GRAY, True, PP_ALIGN.CENTER)
add_text_box(slide, 1.6, 4.2, 2, 0.4, "AUTOMATED", 12, MED_GRAY, True, PP_ALIGN.CENTER)
add_text_box(slide, 10.0, 4.2, 2, 0.4, "HUMAN-GUIDED", 12, MED_GRAY, True, PP_ALIGN.CENTER)

# Competitors
add_text_box(slide, 3.5, 2.8, 2.5, 0.4, "\u25CF  LANDR", 14, ACCENT_BLUE, True)
add_text_box(slide, 3.0, 5.2, 2, 0.4, "\u25CF  Suno", 14, ACCENT_BLUE, True)
add_text_box(slide, 3.2, 5.7, 2, 0.4, "\u25CF  Udio", 14, ACCENT_BLUE, True)
add_text_box(slide, 3.5, 5.0, 2, 0.4, "\u25CF  Moises", 14, ACCENT_BLUE, True)
add_text_box(slide, 9.0, 5.5, 2.5, 0.4, "\u25CF  Songfinch", 14, ACCENT_BLUE, True)
add_text_box(slide, 9.0, 3.2, 2.5, 0.4, "\u25CF  SoundBetter", 14, ACCENT_BLUE, True)

# Major 72 - highlighted
add_shape_with_text(slide, 7.5, 2.6, 3.0, 0.8,
    "\u2605  MAJOR 72", GOLD, BLACK, 18, True)

# ============================================================
# SLIDE 8: COMPETITIVE DETAIL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.5, 11, 0.8, "COMPETITIVE LANDSCAPE", 36, GOLD, True)
add_gold_line(slide, 1.2)

# Comparison table
headers = ["Capability", "Suno", "LANDR", "Moises", "SoundBetter", "Songfinch", "Major 72"]
rows = [
    ["Upload idea \u2192 polished track", "\u2717", "\u2717", "\u2717", "\u2717", "\u2717", "\u2713"],
    ["Named producer taste profiles", "\u2717", "\u2717", "\u2717", "\u2717", "\u2717", "\u2713"],
    ["AI arrangement + human calibration", "\u2717", "\u2717", "\u2717", "\u2717", "\u2717", "\u2713"],
    ["Sheet music / live charts", "\u2717", "\u2717", "\u2717", "\u2717", "\u2717", "\u2713"],
    ["Session musician integration", "\u2717", "\u2717", "\u2717", "\u2713", "\u2717", "\u2713"],
    ["Consumer custom songs", "Partial", "\u2717", "\u2717", "\u2717", "\u2713", "\u2713"],
    ["Full lifecycle: idea \u2192 stage", "\u2717", "\u2717", "\u2717", "\u2717", "\u2717", "\u2713"],
    ["Copyright-defensible output", "\u2717", "N/A", "N/A", "N/A", "\u2713", "\u2713"],
]

cell_w = 1.65
x_start = 0.8
y_start = 1.7

# Headers
for j, header in enumerate(headers):
    w = 2.5 if j == 0 else cell_w
    x = x_start if j == 0 else x_start + 2.5 + (j - 1) * cell_w
    color = GOLD if j == len(headers) - 1 else WHITE
    add_shape_with_text(slide, x, y_start, w, 0.45, header,
                       DARK_GRAY if j < len(headers) - 1 else RGBColor(0x3A, 0x30, 0x10),
                       color, 12, True)

# Data rows
for i, row in enumerate(rows):
    y = y_start + 0.5 + (i * 0.52)
    bg = RGBColor(0x1E, 0x1E, 0x1E) if i % 2 == 0 else RGBColor(0x25, 0x25, 0x25)
    for j, cell in enumerate(row):
        w = 2.5 if j == 0 else cell_w
        x = x_start if j == 0 else x_start + 2.5 + (j - 1) * cell_w
        cell_color = WHITE
        if cell == "\u2713":
            cell_color = GREEN_ACCENT
        elif cell == "\u2717":
            cell_color = RED_ACCENT
        font_s = 11 if j == 0 else 13
        al = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        tf = add_shape_with_text(slide, x, y, w, 0.48, cell, bg, cell_color, font_s, j == 0)
        tf.paragraphs[0].alignment = al

# ============================================================
# SLIDE 9: HEAR IT MVP
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.5, 11, 0.8, "\"HEAR IT\" MVP", 36, GOLD, True)
add_gold_line(slide, 1.2)

add_text_box(slide, 0.8, 1.5, 11, 0.6,
    "The user experience in 5 steps — under 2 minutes total",
    20, WHITE, True)

mvp_steps = [
    ("1", "UPLOAD", "Drag and drop.\nNo forms. No genre\nselection. System\nauto-detects everything."),
    ("2", "CHOOSE", "\"Polish with Eddie's\nSignature Sound\"\nOne button.\nOne decision."),
    ("3", "WAIT", "30-60 seconds.\nShow waveform being\nsculpted. What Eddie's\nmode is doing."),
    ("4", "REVEAL", "Split-screen\nbefore/after.\nOne-click toggle.\nThe \"oh my god\" moment."),
    ("5", "DOWNLOAD", "Polished mix +\nindividual stems +\nsession metadata.\nOne zip. Done."),
]

for i, (num, title, desc) in enumerate(mvp_steps):
    x = 0.4 + (i * 2.55)
    add_shape_with_text(slide, x, 2.4, 2.3, 0.7, f"{num}", GOLD, BLACK, 28, True, MSO_SHAPE.OVAL)
    add_text_box(slide, x, 3.3, 2.3, 0.5, title, 18, GOLD, True, PP_ALIGN.CENTER)
    add_text_box(slide, x, 3.9, 2.3, 2.0, desc, 13, LIGHT_GRAY, False, PP_ALIGN.CENTER)

add_shape_with_text(slide, 2.0, 6.2, 9.0, 0.8,
    "The \"wow\" metric: if >60% of first-time users download the polished version, the engine works.",
    DARK_GRAY, GOLD, 16, True)

# ============================================================
# SLIDE 10: BUSINESS MODEL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.5, 11, 0.8, "BUSINESS MODEL", 36, GOLD, True)
add_gold_line(slide, 1.2)

# Revenue streams
add_text_box(slide, 0.8, 1.5, 5.5, 0.5, "Revenue Streams", 22, WHITE, True)

streams = [
    ("Hear It Subscriptions", "$29-49/mo", "Recurring, predictable"),
    ("Gift It Transactions", "$49-149/song", "High AOV, viral growth"),
    ("Build It (Phase 2)", "$79-149/track", "Per-project upsell"),
    ("Play It (Phase 3)", "$299-999/project", "Premium, high-margin"),
]

for i, (name, price, note) in enumerate(streams):
    y = 2.1 + (i * 0.65)
    add_shape_with_text(slide, 0.8, y, 5.5, 0.55, "", DARK_GRAY)
    add_text_box(slide, 1.0, y + 0.05, 2.2, 0.4, name, 14, WHITE, True)
    add_text_box(slide, 3.2, y + 0.05, 1.3, 0.4, price, 14, GOLD, True, PP_ALIGN.CENTER)
    add_text_box(slide, 4.6, y + 0.05, 1.7, 0.4, note, 12, MED_GRAY, False)

# Year 1 targets
add_text_box(slide, 7.2, 1.5, 5.5, 0.5, "Year 1 Targets", 22, WHITE, True)

targets = [
    ("Hear It subscribers", "1,000"),
    ("Gift It songs sold", "2,000"),
    ("Combined ARR", "~$666K"),
    ("Free-to-paid conversion", ">5%"),
    ("Monthly churn", "<8%"),
    ("Cost per track", "<$0.50"),
]

for i, (metric, target) in enumerate(targets):
    y = 2.1 + (i * 0.55)
    add_shape_with_text(slide, 7.2, y, 5.3, 0.48, "", DARK_GRAY)
    add_text_box(slide, 7.4, y + 0.03, 3.0, 0.4, metric, 13, LIGHT_GRAY, False)
    add_text_box(slide, 10.5, y + 0.03, 1.8, 0.4, target, 14, GOLD, True, PP_ALIGN.RIGHT)

# Future revenue
add_text_box(slide, 0.8, 5.2, 11.5, 0.5,
    "FUTURE REVENUE LAYER: Producer Mode Marketplace",
    18, GOLD, True)
tf = add_text_box(slide, 0.8, 5.7, 11.5, 1.5,
    "Once Eddie's mode is proven, open the platform to other producers. "
    "Producers sell their unique \"taste profiles\" — a new category like sample packs or VST plugins. "
    "Major 72 takes a platform cut on every mode sold and every track processed through it.",
    15, LIGHT_GRAY)

# ============================================================
# SLIDE 11: DIFFERENTIATION / MOAT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.5, 11, 0.8, "DEFENSIBLE ADVANTAGES", 36, GOLD, True)
add_gold_line(slide, 1.2)

moats = [
    ("The Empty Quadrant",
     "No competitor is both professional-grade AND human-guided. Suno is automated + consumer. "
     "LANDR is professional but finishing-only. We're first to this position."),
    ("Taste Can't Be Cloned",
     "Eddie's Creative Mode is built from continuous approve/reject calibration — not a preset. "
     "Competitors can copy features but can't replicate 20 years of aesthetic judgment."),
    ("Copyright by Design",
     "Users upload original material + make creative selections = copyrightable hybrid work. "
     "Unlike prompt-only generators whose output may be public domain."),
    ("The AI Slop Antidote",
     "In a market drowning in generic AI tracks, \"Verified Professional Production\" backed "
     "by a named producer becomes an increasingly valuable quality signal."),
    ("Network Effects (Phase 3+)",
     "Producer Mode Marketplace + Session Musician Network create two-sided platform effects. "
     "More producers = more modes = more users. More musicians = better live pipeline = more projects."),
]

for i, (title, desc) in enumerate(moats):
    y = 1.6 + (i * 1.1)
    add_shape_with_text(slide, 0.8, y, 0.6, 0.5, str(i + 1), GOLD, BLACK, 18, True, MSO_SHAPE.OVAL)
    add_text_box(slide, 1.6, y - 0.05, 3.0, 0.4, title, 17, GOLD, True)
    add_text_box(slide, 1.6, y + 0.35, 10.5, 0.6, desc, 13, LIGHT_GRAY)

# ============================================================
# SLIDE 12: ROADMAP
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.5, 11, 0.8, "ROADMAP", 36, GOLD, True)
add_gold_line(slide, 1.2)

phases = [
    ("PHASE 1", "Months 1-6", "PROVE THE THESIS",
     ["Core engine: stem sep + mix + master",
      "Eddie Signature Creative Mode",
      "\"Hear It\" MVP launch",
      "\"Gift It\" consumer soft-launch",
      "YouTube before/after campaign",
      "Target: 1K subscribers, ~$466K ARR"]),
    ("PHASE 2", "Months 6-12", "EXPAND THE VALUE",
     ["3-4 additional Creative Modes",
      "Session Players (AI arrangement)",
      "\"Build It\" launch",
      "Score generation / live charts",
      "Target: $1-2M ARR"]),
    ("PHASE 3", "Months 12-24", "BUILD THE NETWORK",
     ["\"Play It\" launch",
      "Session musician marketplace",
      "Producer Mode Marketplace (beta)",
      "Sync licensing features",
      "Target: $5M+ ARR"]),
]

for i, (phase, timeline, title, items) in enumerate(phases):
    x = 0.6 + (i * 4.15)
    accent = GOLD if i == 0 else ACCENT_BLUE
    add_shape_with_text(slide, x, 1.5, 3.8, 5.3, "", DARK_GRAY)
    add_text_box(slide, x + 0.2, 1.6, 3.4, 0.4, phase, 14, accent, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, 2.0, 3.4, 0.3, timeline, 12, MED_GRAY, False, PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, 2.4, 3.4, 0.4, title, 16, WHITE, True, PP_ALIGN.CENTER)

    tf = add_text_box(slide, x + 0.3, 3.0, 3.2, 3.5, "", 13, LIGHT_GRAY)
    for item in items:
        add_paragraph(tf, f"\u2022  {item}", 12, LIGHT_GRAY, False, PP_ALIGN.LEFT, Pt(6))

# ============================================================
# SLIDE 13: TEAM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.5, 11, 0.8, "THE TEAM", 36, GOLD, True)
add_gold_line(slide, 1.2)

add_shape_with_text(slide, 1.5, 1.8, 10, 3.5, "", DARK_GRAY)

add_text_box(slide, 2.0, 2.0, 9, 0.5, "Eddie Wilson  —  Founder / Music Director", 24, GOLD, True)
add_text_box(slide, 2.0, 2.6, 9, 2.5,
    "20+ years of studio and live music production experience. Eddie's role is not just advisory — "
    "his creative judgment IS the product. He defines the aesthetic boundaries of each Creative Mode, "
    "evaluates AI output, and calibrates the production parameters that make Major 72's output "
    "indistinguishable from manual professional work.\n\n"
    "As System Architect and Music Director, Eddie brings:\n"
    "\u2022  Deep expertise in mixing, mastering, and arrangement across genres\n"
    "\u2022  A philosophy of \"tools that save time but allow human intervention\"\n"
    "\u2022  The taste and ear that becomes Major 72's core IP\n"
    "\u2022  The vision for a platform where professional judgment scales to infinite tracks",
    14, LIGHT_GRAY)

add_text_box(slide, 1.5, 5.8, 10, 1.0,
    "We are building out the technical team. Key hires needed:\n"
    "CTO / ML Engineering Lead  \u2022  Full-Stack Engineer  \u2022  Audio DSP Specialist",
    14, MED_GRAY, False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 14: THE ASK
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.5, 11, 0.8, "THE ASK", 36, GOLD, True)
add_gold_line(slide, 1.2)

add_text_box(slide, 1.5, 2.0, 10, 0.8,
    "We are raising a seed round to build and launch the\n\"Hear It\" MVP and \"Gift It\" consumer product.",
    22, WHITE, True, PP_ALIGN.CENTER)

uses = [
    ("Core Engine Development", "Stem separation, mixing/mastering pipeline, Creative Mode architecture"),
    ("Eddie Signature Mode", "Calibration tooling, reference track system, approve/reject feedback loop"),
    ("Platform & Infrastructure", "Cloud GPU inference fleet, API, storage, billing, user accounts"),
    ("Go-to-Market", "YouTube content campaign, producer community outreach, Gift It consumer launch"),
    ("Team", "CTO/ML lead, full-stack engineer, audio DSP specialist"),
]

for i, (title, desc) in enumerate(uses):
    y = 3.2 + (i * 0.75)
    add_shape_with_text(slide, 2.0, y, 0.5, 0.5, str(i + 1), GOLD, BLACK, 16, True, MSO_SHAPE.OVAL)
    add_text_box(slide, 2.8, y, 2.5, 0.5, title, 15, GOLD, True)
    add_text_box(slide, 5.3, y, 6.5, 0.5, desc, 13, LIGHT_GRAY)

# ============================================================
# SLIDE 15: CLOSING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 1, 1.5, 11, 1.2, "MAJOR 72 STUDIO", 54, GOLD, True, PP_ALIGN.CENTER)
add_gold_line(slide, 2.8)
add_text_box(slide, 1, 3.2, 11, 1.0,
    "In a world of infinite synthetic content,\n"
    "professional creative judgment is the most valuable currency.",
    24, WHITE, False, PP_ALIGN.CENTER)

add_text_box(slide, 1, 4.8, 11, 0.6,
    "Your Music. Eddie's Ears.",
    30, GOLD, True, PP_ALIGN.CENTER)

add_text_box(slide, 3, 6.0, 7, 0.8,
    "eddie@major72.studio  |  major72.studio",
    16, MED_GRAY, False, PP_ALIGN.CENTER)

# ============================================================
# SAVE
# ============================================================
output_path = os.path.join(os.path.dirname(os.path.dirname(__file__)),
                           "docs", "Major-72-Studio-Pitch-Deck.pptx")
prs.save(output_path)
print(f"Pitch deck saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
